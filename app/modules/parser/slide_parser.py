from __future__ import annotations

import io
import os
from dataclasses import dataclass
from pathlib import Path
from typing import List
import re

from pydantic import ValidationError

from app.schemas.common import BlockType, ParseResponse, SlideBlock, SlidePage
from app.storage import assets
from app.utils.identifiers import new_id
from app.utils.logger import logger

try:  # pragma: no cover - optional dependency guard
    import fitz  # PyMuPDF
except ImportError:  # pragma: no cover - optional dependency guard
    fitz = None
    logger.warning("PyMuPDF unavailable: PDF parsing/image extraction disabled")

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as exc:  # pragma: no cover - dependency optional during tests
    Presentation = None
    MSO_SHAPE_TYPE = None
    logger.warning("python-pptx unavailable: PPTX parsing disabled (%s)", exc)


EMU_PER_INCH = 914400


def _emu_to_points(value: int) -> float:
    return round((value / EMU_PER_INCH) * 72.0, 3)


@dataclass(slots=True)
class ParseResult:
    response: ParseResponse


class SlideParser:
    def parse(self, file_path: Path, file_type: str, session_id: str) -> ParseResponse:
        if file_type == "pdf":
            slides = self._parse_pdf(file_path, session_id)
        elif file_type == "pptx":
            slides = self._parse_pptx(file_path, session_id)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        doc_meta = {"title": file_path.stem, "pages": len(slides)}
        try:
            return ParseResponse(doc_meta=doc_meta, slides=slides)
        except ValidationError as exc:  # pragma: no cover - schema guard
            logger.error("Parse response failed validation: %s", exc)
            raise

    def _parse_pdf(self, file_path: Path, session_id: str) -> List[SlidePage]:
        if fitz is None:
            raise RuntimeError("PyMuPDF 未安装，无法解析 PDF")
        try:
            fitz_doc = fitz.open(str(file_path))
        except Exception as exc:  # pragma: no cover - dependency guard
            logger.error("无法打开 PDF 文件: %s", exc)
            raise

        results: List[SlidePage] = []
        try:
            if fitz_doc.page_count == 0:
                raise ValueError("PDF 未包含任何页面，无法解析")
            for page_index in range(fitz_doc.page_count):
                page = fitz_doc.load_page(page_index)
                blocks: List[SlideBlock] = []
                order = 0
                text_blocks = page.get_text("blocks") or []
                for block in text_blocks:
                    if len(block) < 6:
                        continue
                    x0, y0, x1, y1, text, block_no, block_type = block[:7]
                    if block_type != 0:  # 仅处理文字块
                        continue
                    snippet = (text or "").strip()
                    if not snippet:
                        continue
                    block_kind = BlockType.formula if _likely_formula(snippet) else BlockType.text
                    blocks.append(
                        SlideBlock(
                            id=new_id("b"),
                            type=block_kind,
                            order=order,
                            raw_text=snippet,
                            bbox=[
                                float(x0),
                                float(y0),
                                float(max(x1 - x0, 0.0)),
                                float(max(y1 - y0, 0.0)),
                            ],
                        )
                    )
                    order += 1
                if not blocks:
                    fallback_text = (page.get_text("text") or "").strip()
                    if fallback_text:
                        block_type = BlockType.formula if _likely_formula(fallback_text) else BlockType.text
                        blocks.append(
                            SlideBlock(
                                id=new_id("b"),
                                type=block_type,
                                order=order,
                                raw_text=fallback_text,
                                bbox=[float(page.rect.width), float(page.rect.height), 0.0, 0.0],
                            )
                        )
                        order += 1

                image_blocks = self._extract_pdf_images(
                    fitz_doc,
                    page_index + 1,
                    session_id,
                    float(page.rect.width),
                    float(page.rect.height),
                )
                for image_block in image_blocks:
                    image_block.order = order
                    blocks.append(image_block)
                    order += 1

                results.append(SlidePage(page_no=page_index + 1, blocks=blocks))
        finally:
            fitz_doc.close()
        return results

    def _extract_pdf_images(
        self,
        fitz_doc,
        page_no: int,
        session_id: str,
        page_width: float,
        page_height: float,
    ) -> List[SlideBlock]:
        if fitz_doc is None:
            return []
        try:
            page = fitz_doc.load_page(page_no - 1)
        except Exception as exc:  # pragma: no cover - defensive
            logger.warning("加载 PDF 页面失败，无法抽取图片: page=%s error=%s", page_no, exc)
            return []
        dict_blocks = page.get_text("dict").get("blocks", [])
        entries: List[tuple[tuple[float, float, float, float], int]] = []
        seen_xrefs = set()
        for block in dict_blocks:
            if block.get("type") != 1:
                continue
            xref = block.get("image")
            bbox = block.get("bbox")
            if xref is None or bbox is None:
                continue
            if xref in seen_xrefs:
                continue
            seen_xrefs.add(xref)
            entries.append((tuple(float(value) for value in bbox), xref))
        if not entries:
            # Fallback：采用 get_images 列表，但无 bbox 信息，退化为整页
            for image in page.get_images(full=True):
                xref = image[0]
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)
                entries.append(((0.0, 0.0, float(page_width), float(page_height)), xref))
        if not entries:
            return []
        entries.sort(key=lambda item: (item[0][1], item[0][0]))
        blocks: List[SlideBlock] = []
        for bbox, xref in entries:
            try:
                base_image = fitz_doc.extract_image(xref)
            except Exception as exc:  # pragma: no cover - defensive
                logger.debug("提取 PDF 图片失败: page=%s xref=%s error=%s", page_no, xref, exc)
                continue
            image_bytes = base_image.get("image")
            if not image_bytes:
                continue
            ext = base_image.get("ext", "png") or "png"
            rel_name = f"{new_id('img')}_p{page_no}_xref{xref}.{ext}"
            asset_uri = assets.write_asset(session_id, rel_name, image_bytes)
            blocks.append(
                SlideBlock(
                    id=new_id("b"),
                    type=BlockType.image,
                    order=0,  # 由调用方重排
                    bbox=[float(coord) for coord in bbox],
                    asset_uri=asset_uri,
                )
            )
        return blocks

    def _parse_pptx(self, file_path: Path, session_id: str) -> List[SlidePage]:
        if Presentation is None:
            raise RuntimeError(
                "python-pptx is required to parse PPTX files. Install python-pptx."
            )
        pres = Presentation(str(file_path))
        slides: List[SlidePage] = []
        for idx, slide in enumerate(pres.slides, start=1):
            blocks: List[SlideBlock] = []
            order = 0
            for shape in slide.shapes:
                bbox = [
                    _emu_to_points(shape.left),
                    _emu_to_points(shape.top),
                    _emu_to_points(shape.width),
                    _emu_to_points(shape.height),
                ]
                if shape.has_text_frame and shape.text.strip():
                    text = shape.text.strip()
                    if order == 0:
                        block_type = BlockType.title
                    elif _likely_formula(text):
                        block_type = BlockType.formula
                    else:
                        block_type = BlockType.text
                    blocks.append(
                        SlideBlock(
                            id=new_id("b"),
                            type=block_type,
                            order=order,
                            raw_text=text,
                            bbox=bbox,
                        )
                    )
                    order += 1
                elif MSO_SHAPE_TYPE and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    stream = io.BytesIO(image.blob)
                    rel_name = f"{new_id('img')}.{image.ext or 'png'}"
                    asset_uri = assets.write_asset(session_id, rel_name, stream.read())
                    blocks.append(
                        SlideBlock(
                            id=new_id("b"),
                            type=BlockType.image,
                            order=order,
                            bbox=bbox,
                            asset_uri=asset_uri,
                        )
                    )
                    order += 1
            slides.append(SlidePage(page_no=idx, blocks=blocks))
        return slides


FORMULA_PATTERN = re.compile(r"(\\[a-zA-Z]+|[=±×÷∑∫√^_])")


def _likely_formula(text: str) -> bool:
    return bool(FORMULA_PATTERN.search(text))
